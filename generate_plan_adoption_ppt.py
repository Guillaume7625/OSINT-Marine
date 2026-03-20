#!/usr/bin/env python3
"""
Génère un PowerPoint (.pptx) à partir d'un plan de slides (surtitre/titre/bullets/punchline).

Dépendance :
  pip install python-pptx

Exemples :
  python3 generate_plan_adoption_ppt.py
  python3 generate_plan_adoption_ppt.py --output plan_adoption_ia.pptx
  python3 generate_plan_adoption_ppt.py --template template_corporate.pptx --output deck.pptx
"""

from __future__ import annotations

import argparse
import re
from dataclasses import dataclass
from typing import Iterable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


def _split_md_bold(text: str) -> list[tuple[str, bool]]:
    """Split text into (segment, is_bold) based on **...** markers."""
    parts: list[tuple[str, bool]] = []
    pos = 0
    for m in _BOLD_RE.finditer(text):
        if m.start() > pos:
            parts.append((text[pos : m.start()], False))
        parts.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        parts.append((text[pos:], False))
    return parts


def _set_paragraph_runs(paragraph, text: str, *, font_name: str, font_size_pt: int, color: RGBColor):
    paragraph.text = ""
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size_pt)
    paragraph.font.color.rgb = color
    for seg, is_bold in _split_md_bold(text):
        if not seg:
            continue
        run = paragraph.add_run()
        run.text = seg
        run.font.bold = is_bold


def _set_shape_text(shape, text: str, *, font_name: str, font_size_pt: int, color: RGBColor, bold: bool = False, italic: bool = False):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size_pt)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic


@dataclass(frozen=True)
class SlideSpec:
    surtitre: str
    titre: str
    bullets: list[str]
    punchline: str
    note: Optional[str] = None
    callout: Optional[str] = None


def _make_presentation(*, template: Optional[str], widescreen: bool) -> Presentation:
    prs = Presentation(template) if template else Presentation()
    if not template and widescreen:
        # 16:9 -> width = 7.5 * 16/9
        prs.slide_height = Inches(7.5)
        prs.slide_width = Inches(7.5 * 16 / 9)
    return prs


def _add_slide(prs: Presentation, spec: SlideSpec):
    # Layout "Title and Content" pour hériter du style de bullets du master.
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    font_name = "Calibri"
    color_title = RGBColor(0x11, 0x11, 0x11)
    color_body = RGBColor(0x11, 0x11, 0x11)
    color_surtitre = RGBColor(0x66, 0x66, 0x66)
    color_accent = RGBColor(0x1F, 0x4E, 0x79)  # bleu sobre

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin_x = Inches(0.7)
    top_margin = Inches(0.35)
    gap = Inches(0.12)

    surtitre_h = Inches(0.30)
    title_h = Inches(1.05)
    punch_h = Inches(0.55)
    punch_top = slide_h - Inches(0.35) - punch_h

    reserved = []
    if spec.callout:
        reserved.append(("callout", Inches(0.80)))
    if spec.note:
        reserved.append(("note", Inches(0.35)))

    reserved_total = sum(h for _, h in reserved) + (gap * (len(reserved) if reserved else 0))

    content_w = slide_w - (2 * margin_x)
    surtitre_top = top_margin
    title_top = surtitre_top + surtitre_h + Inches(0.05)
    body_top = title_top + title_h + Inches(0.05)
    body_h = max(Inches(1.0), punch_top - body_top - gap - reserved_total)

    # Surtitre (textbox)
    surtitre_box = slide.shapes.add_textbox(margin_x, surtitre_top, content_w, surtitre_h)
    _set_shape_text(
        surtitre_box,
        spec.surtitre,
        font_name=font_name,
        font_size_pt=14,
        color=color_surtitre,
        bold=True,
    )

    # Titre (placeholder)
    title_shape = slide.shapes.title
    title_shape.left = margin_x
    title_shape.top = title_top
    title_shape.width = content_w
    title_shape.height = title_h
    _set_shape_text(
        title_shape,
        spec.titre,
        font_name=font_name,
        font_size_pt=34,
        color=color_title,
        bold=True,
    )

    # Bullets (placeholder)
    body_shape = slide.placeholders[1]
    body_shape.left = margin_x
    body_shape.top = body_top
    body_shape.width = content_w
    body_shape.height = body_h
    tf = body_shape.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, bullet in enumerate(spec.bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        _set_paragraph_runs(
            p,
            bullet,
            font_name=font_name,
            font_size_pt=20,
            color=color_body,
        )

    # Zones réservées (callout/note) entre bullets et punchline
    y = body_top + body_h + gap
    for kind, h in reserved:
        if kind == "callout" and spec.callout:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin_x, y, content_w, h)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
            box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            tfc = box.text_frame
            tfc.clear()
            tfc.word_wrap = True
            p = tfc.paragraphs[0]
            _set_paragraph_runs(
                p,
                f"**Encadré :** {spec.callout}",
                font_name=font_name,
                font_size_pt=14,
                color=color_body,
            )

        if kind == "note" and spec.note:
            note_box = slide.shapes.add_textbox(margin_x, y, content_w, h)
            _set_shape_text(
                note_box,
                spec.note,
                font_name=font_name,
                font_size_pt=12,
                color=color_surtitre,
                italic=True,
            )
        y = y + h + gap

    # Punchline (textbox)
    punch_box = slide.shapes.add_textbox(margin_x, punch_top, content_w, punch_h)
    _set_shape_text(
        punch_box,
        spec.punchline,
        font_name=font_name,
        font_size_pt=18,
        color=color_accent,
        italic=True,
    )


def build_deck(*, template: Optional[str], widescreen: bool) -> Presentation:
    prs = _make_presentation(template=template, widescreen=widescreen)

    slides: list[SlideSpec] = [
        SlideSpec(
            surtitre="Le constat",
            titre="Les outils sont là — les pratiques ne suivent pas",
            bullets=[
                "Les outils existent : GENIAL, Qlik, UiPath — les 5 chantiers sont lancés",
                "Les pratiques restent inégales d'un bureau à l'autre",
                "Dès que la charge monte, les nouveaux usages disparaissent",
                "**74 %** des organisations qui ont les bons outils n'en tirent aucune valeur — le problème est toujours l'adoption, pas la technologie (BCG, 2024)",
            ],
            punchline="L'outil sans l'usage reste une vitrine — l'usage sans l'outil reste un discours",
        ),
        SlideSpec(
            surtitre="La réponse",
            titre="Embarquer trois niveaux en même temps",
            bullets=[
                "Les autorités utilisent l'IA sur un cas concret → elles savent quoi exiger",
                "Les chefs de bureau transforment un processus → le gain est mesuré",
                "Les collaborateurs pratiquent sur leurs tâches dès J+7 → l'usage s'ancre ou non",
            ],
            punchline="Le haut lance, le terrain prouve, la boucle se ferme",
        ),
        SlideSpec(
            surtitre="Les autorités — pourquoi et comment",
            titre="Pour exiger juste, il faut avoir vu",
            bullets=[
                "L'IA réduit des jours d'analyse à quelques heures",
                "Ce gain reste abstrait tant qu'on ne l'a pas constaté soi-même",
                "Sans avoir pratiqué, impossible d'exiger juste ni de juger un résultat",
                "L'amiral Vandier a fait de la formation IA une priorité pour ses officiers au Commandement suprême allié Transformation",
            ],
            callout="Des experts IA accompagnent chaque autorité individuellement, sur ses sujets, à son rythme",
            punchline="L'autorité qui a pratiqué arbitre mieux — et son exigence devient crédible",
        ),
        SlideSpec(
            surtitre="Les chefs de bureau — pourquoi et comment",
            titre="C'est le chef de bureau qui décide si l'IA entre dans les habitudes",
            bullets=[
                "C'est lui qui lance le mouvement et protège le temps pour pratiquer",
                "Pour s'engager, il doit voir un gain concret sur son propre processus",
                "Kick-off tous chefs de bureau (½ journée) : cadre et attendus",
                "Kick-off 3 bureaux pilotes (½ journée) : choix du processus cible",
                "Les pilotes sont accompagnés 2h par semaine pendant 6 semaines",
                "Les autres bureaux récupèrent les méthodes et résultats des pilotes",
            ],
            punchline="Quand le chef voit le gain, il protège le temps — et l'usage tient",
        ),
        SlideSpec(
            surtitre="Les collaborateurs — ce qui bloque",
            titre="Ils ont essayé — et ça les a déçus",
            bullets=[
                "La plupart ont utilisé l'IA comme un moteur de recherche",
                "Question vague, réponse médiocre, conclusion : l'outil ne marche pas",
                "Le problème n'est pas l'outil — personne ne leur a montré comment s'en servir",
                "Sans formation, chaque mauvais essai renforce le rejet",
            ],
            punchline="Le premier essai sans formation est souvent le dernier",
        ),
        SlideSpec(
            surtitre="Les collaborateurs — la preuve que ça se résout",
            titre="Avec la bonne formation, les résultats sont mesurés",
            bullets=[
                "**+43 %** de performance pour ceux en difficulté",
                "**+17 %** pour ceux déjà bons",
                "**+12 %** de tâches accomplies en plus",
                "**+25 %** de rapidité",
            ],
            note="(Dell'Acqua et al., Harvard Business School, 2023 — 758 consultants BCG)",
            punchline="La formation fait toute la différence",
        ),
        SlideSpec(
            surtitre="Les collaborateurs — ce qu'on leur propose",
            titre="Une demi-journée pour corriger le premier essai, 6 semaines pour ancrer",
            bullets=[
                "Une demi-journée par collaborateur sur son cas réel — avec GENIAL, les mini-apps ou UiPath",
                "L'objectif : un cas réussi dès la première demi-journée",
                "Tous peuvent ensuite suivre une formation de 6 semaines sur Mistral AI (données non sensibles)",
                "Résultat collectif : des prompts testés et validés, utilisables dès le lendemain",
            ],
            punchline="L'agent formé utilise l'outil pour de vrai — pas en démonstration",
        ),
        SlideSpec(
            surtitre="Le calendrier",
            titre="Tout tient en 3 mois",
            bullets=[
                "**Semaine 1** — Accompagnement des autorités + kick-off tous chefs de bureau (½ j) + kick-off bureaux pilotes (½ j)",
                "**Semaine 2** — Demi-journées collaborateurs sur cas réels + ouverture formation Mistral AI",
                "**S2 à S7** — Rendez-vous hebdomadaires bureaux pilotes + formation continue collaborateurs",
                "**Mars** — Point de mi-parcours : ce qui ne produit pas d'usage réel est ajusté ou arrêté",
                "**Juin** — Bilan : adoption 70 %, pilotes 2/3, conformité 100 %",
            ],
            punchline="Dans 3 mois, l'EMM sait si SIGNAL change le quotidien — avec des faits",
        ),
    ]

    for spec in slides:
        _add_slide(prs, spec)

    return prs


def parse_args(argv: Optional[Iterable[str]] = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Génère le PPT 'plan d’adoption IA' (8 slides).")
    parser.add_argument("--output", default="plan_adoption_ia.pptx", help="Chemin du .pptx à générer.")
    parser.add_argument("--template", default=None, help="PPTX template (thème corporate) à utiliser.")
    parser.add_argument(
        "--no-widescreen",
        action="store_true",
        help="Désactive le format 16:9 (utile si vous voulez garder le 4:3 par défaut).",
    )
    return parser.parse_args(argv)


def main(argv: Optional[Iterable[str]] = None) -> int:
    args = parse_args(argv)
    prs = build_deck(template=args.template, widescreen=not args.no_widescreen)
    prs.save(args.output)
    print(f"OK: {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

